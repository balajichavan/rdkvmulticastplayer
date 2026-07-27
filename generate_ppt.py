#!/usr/bin/env python3
"""Generate a PowerPoint deck for IP Multicast Playback in an RDKV Player App."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- Theme colors ----
DARK = RGBColor(0x0B, 0x1F, 0x3A)
ACCENT = RGBColor(0x1E, 0x88, 0xE5)
ACCENT2 = RGBColor(0x00, 0xAC, 0xC1)
LIGHT = RGBColor(0xF4, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x5A, 0x6B, 0x7B)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5A, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(BLANK)


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0, round_=True):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def txt(slide, x, y, w, h, text, size=18, color=DARK, bold=False,
        align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font="Segoe UI", italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = font
    return tb


def shape_text(shp, text, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER):
    tf = shp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    lines = text.split("\n")
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = ln
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Segoe UI"


def arrow(slide, x, y, w, h, color=ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    a = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def down_arrow(slide, x, y, w, h, color=ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    a = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    a.fill.solid()
    a.fill.fore_color.rgb = color
    a.line.fill.background()
    a.shadow.inherit = False
    return a


def header(slide, title, subtitle=None):
    bar = box(slide, 0, 0, SW, Inches(1.05), fill=DARK, round_=False)
    box(slide, 0, Inches(1.05), SW, Inches(0.08), fill=ACCENT, round_=False)
    txt(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.6),
        title, size=28, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        txt(slide, Inches(0.52), Inches(0.66), Inches(12), Inches(0.35),
            subtitle, size=13, color=RGBColor(0xB8, 0xC7, 0xD9))


# ============ Slide 1: Title ============
s = add_slide()
bg(s, DARK)
box(s, 0, Inches(2.4), SW, Inches(0.09), fill=ACCENT, round_=False)
box(s, 0, Inches(4.55), SW, Inches(0.09), fill=ACCENT2, round_=False)
txt(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.4),
    "IP Multicast Playback in an\nRDKV Player App", size=44, color=WHITE, bold=True)
txt(s, Inches(0.85), Inches(4.75), Inches(11.5), Inches(0.6),
    "Suggested Architecture & Integration Options", size=22, color=ACCENT2)
txt(s, Inches(0.85), Inches(6.6), Inches(11.5), Inches(0.4),
    "Player UI App  •  Thunder / WPEFramework  •  AAMP  •  GStreamer", size=13,
    color=RGBColor(0x9F, 0xB2, 0xC8))

# ============ Slide 2: Objective & Principles ============
s = add_slide()
bg(s, LIGHT)
header(s, "Objective & Key Principles")
txt(s, Inches(0.6), Inches(1.4), Inches(12), Inches(0.7),
    "Play an IP multicast stream (MPEG-TS over UDP/RTP) inside an RDKV player UI app "
    "using the platform's hardware demux/decode pipeline, controlled from the UI via "
    "Thunder (WPEFramework) JSON-RPC.", size=16, color=GREY)
principles = [
    ("Do NOT use a generic element", "Never feed a raw multicast URL to an HTML <video> element."),
    ("Route through RDK media stack", "Use hardware decode via westeros / brcm sinks."),
    ("Platform manages IGMP", "IGMP join/leave and decoder resources handled by platform."),
    ("Keep the UI portable", "UI only issues JSON-RPC control calls — no pipeline logic."),
]
y = Inches(2.5)
for i, (t, d) in enumerate(principles):
    row = y + Inches(i * 1.05)
    box(s, Inches(0.6), row, Inches(0.12), Inches(0.85), fill=ACCENT, round_=False)
    txt(s, Inches(0.9), row, Inches(4.6), Inches(0.85), t, size=17, color=DARK,
        bold=True, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, Inches(5.7), row, Inches(7.0), Inches(0.85), d, size=15, color=GREY,
        anchor=MSO_ANCHOR.MIDDLE)

# ============ Slide 3: Option 1 AAMP MABR ============
s = add_slide()
bg(s, LIGHT)
header(s, "Option 1 — AAMP with Multicast / MABR", "Recommended for production RDKV apps")

# pill "RECOMMENDED"
pill = box(s, Inches(9.9), Inches(0.28), Inches(2.9), Inches(0.5), fill=GREEN)
shape_text(pill, "RECOMMENDED", size=13, color=WHITE, bold=True)

lane_y = Inches(1.55)
lane_h = Inches(1.15)
w = Inches(2.5)
gap = Inches(0.35)
x0 = Inches(0.55)


def flow_box(slide, idx, title, sub, color):
    x = x0 + idx * (w + gap + Inches(0.35))
    b = box(slide, x, lane_y, w, lane_h, fill=color)
    tf = b.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title
    r.font.size = Pt(13); r.font.bold = True; r.font.color.rgb = WHITE
    r.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sub
    r2.font.size = Pt(9.5); r2.font.color.rgb = RGBColor(0xE3, 0xEE, 0xFB)
    r2.font.name = "Segoe UI"
    return x


stages = [
    ("UI App", "Lightning / React JS", ACCENT),
    ("Thunder Plugin", "org.rdk.MediaPlayer", RGBColor(0x37, 0x47, 0x5A)),
    ("AAMP Core", "MABR + TS Demux + DRM", ACCENT2),
    ("GStreamer", "brcmdecoder → westerossink", RGBColor(0x37, 0x47, 0x5A)),
]
positions = []
for i, (t, sd, c) in enumerate(stages):
    positions.append(flow_box(s, i, t, sd, c))
for i in range(len(stages) - 1):
    ax = positions[i] + w
    arrow(s, ax + Inches(0.02), lane_y + Inches(0.42), Inches(0.33), Inches(0.3), ACCENT)

# source + display
src = box(s, x0, Inches(3.25), w, Inches(0.85), fill=ORANGE)
shape_text(src, "IP Multicast Source\nMPEG-TS UDP/RTP", size=12)
down_arrow(s, positions[2] + Inches(0.95), Inches(2.75), Inches(0.45), Inches(0.4), ACCENT)
disp = box(s, positions[3], Inches(3.25), w, Inches(0.85), fill=DARK)
shape_text(disp, "Display / HDMI", size=13)
# IGMP arrow src->AAMP
arrow(s, x0 + w + Inches(0.1), Inches(3.5), Inches(1.0), Inches(0.3), ORANGE)
txt(s, x0 + w + Inches(0.05), Inches(3.15), Inches(1.5), Inches(0.3), "IGMP join",
    size=9, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

feats = [
    "Standard RDK-V player — built-in Multicast/MABR helper (mcast-to-unicast proxy)",
    "Supports ABR, DRM/CAS via OCDM, and hardware decode",
    "UI controls via org.rdk.MediaPlayer JSON-RPC",
    "Locator example:  mcast://<group-ip>:<port>",
]
yb = Inches(4.55)
for i, f in enumerate(feats):
    box(s, Inches(0.6), yb + Inches(i * 0.6) + Inches(0.05), Inches(0.14), Inches(0.14),
        fill=GREEN, round_=False)
    txt(s, Inches(0.95), yb + Inches(i * 0.6) - Inches(0.05), Inches(11.6), Inches(0.5),
        f, size=14, color=GREY)

# ============ Slide 4: Option 2 Native GStreamer ============
s = add_slide()
bg(s, LIGHT)
header(s, "Option 2 — Native GStreamer via Custom Thunder Plugin",
       "For simple, unencrypted CBR multicast")

lane_y = Inches(1.7)
stages2 = [
    ("UI App", "Lightning JS", ACCENT),
    ("Custom Plugin", "play / stop / tune", RGBColor(0x37, 0x47, 0x5A)),
    ("GStreamer", "udpsrc/rtpbin → tsdemux", ACCENT2),
    ("HW Decode", "brcmdecoder → westerossink", RGBColor(0x37, 0x47, 0x5A)),
]
positions = []
for i, (t, sd, c) in enumerate(stages2):
    x = x0 + i * (w + gap + Inches(0.35))
    b = box(s, x, lane_y, w, lane_h, fill=c)
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = t; r.font.size = Pt(13); r.font.bold = True
    r.font.color.rgb = WHITE; r.font.name = "Segoe UI"
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = sd; r2.font.size = Pt(9.5)
    r2.font.color.rgb = RGBColor(0xE3, 0xEE, 0xFB); r2.font.name = "Segoe UI"
    positions.append(x)
for i in range(len(stages2) - 1):
    ax = positions[i] + w
    arrow(s, ax + Inches(0.02), lane_y + Inches(0.42), Inches(0.33), Inches(0.3), ACCENT)

src = box(s, x0, Inches(3.35), w, Inches(0.85), fill=ORANGE)
shape_text(src, "IP Multicast Source\nMPEG-TS UDP/RTP", size=12)
down_arrow(s, positions[2] + Inches(0.95), Inches(2.9), Inches(0.45), Inches(0.4), ACCENT)
disp = box(s, positions[3], Inches(3.35), w, Inches(0.85), fill=DARK)
shape_text(disp, "Display / HDMI", size=13)

pts = [
    "Pipeline: udpsrc/rtpbin → (rtpmp2tdepay) → tsdemux → brcmvideodecoder → westerossink",
    "No ABR; DRM handled manually",
    "IGMP join performed on device by udpsrc",
    "Wrap pipeline in a custom Thunder plugin exposing play / stop / tune",
]
yb = Inches(4.65)
for i, f in enumerate(pts):
    box(s, Inches(0.6), yb + Inches(i * 0.58) + Inches(0.05), Inches(0.14), Inches(0.14),
        fill=ACCENT2, round_=False)
    txt(s, Inches(0.95), yb + Inches(i * 0.58) - Inches(0.05), Inches(11.6), Inches(0.5),
        f, size=14, color=GREY)

# ============ Slide 5: Option 3 Edge M2U ============
s = add_slide()
bg(s, LIGHT)
header(s, "Option 3 — MSO Edge Multicast-to-Unicast (ABR)",
       "Best when the network already has an edge packager")

# Headend group
box(s, Inches(0.55), Inches(1.55), Inches(4.0), Inches(2.6), fill=RGBColor(0xE7, 0xEE, 0xF6),
    line=ACCENT, line_w=1.5)
txt(s, Inches(0.7), Inches(1.65), Inches(3.7), Inches(0.35), "MSO / Network Edge",
    size=13, color=DARK, bold=True)
b1 = box(s, Inches(0.85), Inches(2.2), Inches(3.4), Inches(0.8), fill=ORANGE)
shape_text(b1, "IP Multicast\nMPEG-TS", size=12)
down_arrow(s, Inches(2.35), Inches(3.02), Inches(0.4), Inches(0.35), ORANGE)
b2 = box(s, Inches(0.85), Inches(3.35), Inches(3.4), Inches(0.7), fill=ACCENT2)
shape_text(b2, "M2U / JIT ABR Packager", size=12)

arrow(s, Inches(4.6), Inches(2.6), Inches(0.7), Inches(0.4), ACCENT)
txt(s, Inches(4.35), Inches(2.2), Inches(1.4), Inches(0.35), "HLS/DASH\nunicast", size=9,
    color=GREY, bold=True, align=PP_ALIGN.CENTER)

# Device group
box(s, Inches(5.5), Inches(1.55), Inches(7.2), Inches(2.6), fill=RGBColor(0xE7, 0xEE, 0xF6),
    line=ACCENT2, line_w=1.5)
txt(s, Inches(5.65), Inches(1.65), Inches(5), Inches(0.35), "RDKV Device", size=13,
    color=DARK, bold=True)
d1 = box(s, Inches(5.75), Inches(2.15), Inches(2.0), Inches(0.75), fill=ACCENT)
shape_text(d1, "UI App", size=12)
d2 = box(s, Inches(8.0), Inches(2.15), Inches(2.1), Inches(0.75), fill=RGBColor(0x37, 0x47, 0x5A))
shape_text(d2, "AAMP Plugin", size=12)
d3 = box(s, Inches(10.35), Inches(2.15), Inches(2.1), Inches(0.75), fill=ACCENT2)
shape_text(d3, "AAMP ABR", size=12)
arrow(s, Inches(7.78), Inches(2.4), Inches(0.2), Inches(0.28), ACCENT)
arrow(s, Inches(10.12), Inches(2.4), Inches(0.2), Inches(0.28), ACCENT)
d4 = box(s, Inches(8.0), Inches(3.25), Inches(4.45), Inches(0.7), fill=DARK)
shape_text(d4, "GStreamer → westerossink (HW decode) → Display", size=12)
down_arrow(s, Inches(11.2), Inches(2.92), Inches(0.4), Inches(0.35), ACCENT)

pts = [
    "Network edge converts multicast to HLS/DASH unicast",
    "Device performs no IGMP join — standard AAMP ABR fetch",
    "Highest portability across SoCs; lowest device-side effort",
]
yb = Inches(4.6)
for i, f in enumerate(pts):
    box(s, Inches(0.6), yb + Inches(i * 0.58) + Inches(0.05), Inches(0.14), Inches(0.14),
        fill=ACCENT, round_=False)
    txt(s, Inches(0.95), yb + Inches(i * 0.58) - Inches(0.05), Inches(11.6), Inches(0.5),
        f, size=14, color=GREY)

# ============ Slide 6: Comparison table ============
s = add_slide()
bg(s, LIGHT)
header(s, "Options Comparison")

rows = [
    ["Aspect", "Option 1: AAMP MABR", "Option 2: Native GStreamer", "Option 3: Edge M2U ABR"],
    ["Best for", "Standard RDK multicast IPTV", "Simple unencrypted CBR", "Networks with edge packagers"],
    ["IGMP join", "On device (AAMP helper)", "On device (udpsrc)", "At network edge"],
    ["ABR support", "Yes (MABR)", "No", "Yes"],
    ["DRM / CAS", "Yes (OCDM)", "Manual", "Yes"],
    ["Portability (SoC)", "High", "Medium", "High"],
    ["Dev effort", "Low–Medium", "Medium–High", "Low (device)"],
]
tx, ty = Inches(0.55), Inches(1.5)
col_w = [Inches(2.6), Inches(3.4), Inches(3.4), Inches(3.4)]
row_h = Inches(0.72)
for ri, row in enumerate(rows):
    cx = tx
    for ci, cell in enumerate(row):
        if ri == 0:
            fill = DARK
        elif ci == 0:
            fill = RGBColor(0xD8, 0xE2, 0xEE)
        else:
            fill = WHITE if ri % 2 else RGBColor(0xEE, 0xF3, 0xF9)
        c = box(s, cx, ty + ri * row_h, col_w[ci], row_h, fill=fill,
                line=RGBColor(0xCF, 0xD9, 0xE4), line_w=0.75, round_=False)
        tf = c.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = cell
        r.font.size = Pt(12.5)
        r.font.bold = (ri == 0 or ci == 0)
        r.font.color.rgb = WHITE if ri == 0 else DARK
        r.font.name = "Segoe UI"
        cx += col_w[ci]

# ============ Slide 7: Recommendation & Next Steps ============
s = add_slide()
bg(s, LIGHT)
header(s, "Recommendation & Next Steps")

# recommendation cards
cards = [
    ("PRIMARY", "Option 1 — AAMP with MABR", "Production RDKV apps", GREEN),
    ("IF EDGE EXISTS", "Option 3 — Edge M2U ABR", "Network has a packager", ACCENT),
    ("LIGHTWEIGHT", "Option 2 — Native GStreamer", "Unencrypted CBR trials", ACCENT2),
]
cw = Inches(3.9)
cx = Inches(0.6)
for i, (tag, title, sub, col) in enumerate(cards):
    x = cx + i * (cw + Inches(0.35))
    box(s, x, Inches(1.5), cw, Inches(1.9), fill=WHITE, line=col, line_w=2)
    box(s, x, Inches(1.5), cw, Inches(0.45), fill=col, round_=False)
    txt(s, x, Inches(1.52), cw, Inches(0.4), tag, size=12, color=WHITE, bold=True,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    txt(s, x + Inches(0.2), Inches(2.15), cw - Inches(0.4), Inches(0.7), title, size=16,
        color=DARK, bold=True, align=PP_ALIGN.CENTER)
    txt(s, x + Inches(0.2), Inches(2.85), cw - Inches(0.4), Inches(0.4), sub, size=13,
        color=GREY, align=PP_ALIGN.CENTER)

txt(s, Inches(0.6), Inches(3.75), Inches(12), Inches(0.4), "Next Steps", size=20,
    color=DARK, bold=True)
steps = [
    "Confirm current player engine (AAMP, custom Thunder plugin, or Lightning template).",
    "Define multicast source format (UDP vs RTP, CBR vs ABR, encrypted vs clear).",
    "Implement Option 1 JSON-RPC integration and AAMP locator wiring.",
    "Validate IGMP join, decoder resource handling, and DRM path.",
]
for i, st in enumerate(steps):
    yv = Inches(4.4) + Inches(i * 0.62)
    n = box(s, Inches(0.6), yv, Inches(0.5), Inches(0.5), fill=ACCENT)
    shape_text(n, str(i + 1), size=16, color=WHITE, bold=True)
    txt(s, Inches(1.3), yv, Inches(11.4), Inches(0.5), st, size=15, color=GREY,
        anchor=MSO_ANCHOR.MIDDLE)

out = "/home/buv115/proposals/playeruiapp/RDKV_Multicast_Playback.pptx"
prs.save(out)
print("Saved:", out)
